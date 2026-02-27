"""
AI Studio - Python Sidecar Server
=================================

FastAPI-based server providing multi-provider LLM chat capabilities.
Connects to local LLMs (Ollama) and cloud providers (Anthropic, OpenAI).
MCP-native tool system with built-in + external tools.

Run modes:
- Development: `python server.py`
- Docker: `docker compose up`
"""

import os
import time
import uuid
import asyncio
from contextlib import asynccontextmanager
from typing import Optional

from fastapi import FastAPI, HTTPException, Request, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel

from agent.chat import ChatService
from agent.embedding import create_embedding_client
from agent.providers import (
    OllamaProvider,
    AnthropicProvider,
    OpenAIProvider,
    GoogleProvider,
    AzureOpenAIProvider,
    LocalOpenAIProvider,
    AgentProvider,
    Message,
)
from agent.mcp import ToolRegistry, McpClientManager, register_builtin_tools
from agent.events import EventBus
from agent.tools import ShellTool, FilesystemTool


# ============================================================================
# Request/Response Models
# ============================================================================

class ChatRequest(BaseModel):
    """Chat completion request"""
    conversation_id: Optional[str] = None
    message: str
    provider: Optional[str] = None
    model: Optional[str] = None
    temperature: float = 0.7
    system_prompt: Optional[str] = None
    # Per-request provider config (passed from Rust, loaded from SQLite settings)
    api_key: Optional[str] = None
    base_url: Optional[str] = None
    extra_config: Optional[dict] = None
    # Tool support
    tools_enabled: bool = False
    # Message history from Rust (SQLite source of truth).
    # When provided, replaces sidecar's in-memory history for this conversation.
    history: Optional[list[dict]] = None


class ChatMessageRequest(BaseModel):
    """Direct message request (no conversation)"""
    messages: list[dict]
    provider: Optional[str] = "ollama"
    model: Optional[str] = None
    temperature: float = 0.7
    # Per-request provider config (for workflow LLM nodes)
    api_key: Optional[str] = None
    base_url: Optional[str] = None
    extra_config: Optional[dict] = None
    system_prompt: Optional[str] = None
    images: Optional[list[dict]] = None  # [{"data": "base64...", "mime_type": "image/png"}]
    # Session mode: accumulate history across calls with the same conversation_id
    conversation_id: Optional[str] = None
    max_history: Optional[int] = 20


class ChatResponse(BaseModel):
    """Chat completion response"""
    conversation_id: str
    content: str
    model: str
    provider: str
    usage: Optional[dict] = None
    tool_calls: Optional[list[dict]] = None


class McpConnectRequest(BaseModel):
    """Request to connect to an MCP server"""
    name: str
    transport: str = "stdio"
    command: Optional[str] = None
    args: list[str] = []
    url: Optional[str] = None
    env: dict[str, str] = {}


class McpDisconnectRequest(BaseModel):
    """Request to disconnect from an MCP server"""
    name: str


# ============================================================================
# Application Setup
# ============================================================================

chat_service: ChatService = None
tool_registry: ToolRegistry = None
mcp_client: McpClientManager = None
event_bus: EventBus = None


@asynccontextmanager
async def lifespan(app: FastAPI):
    """Initialize services on startup"""
    global chat_service, tool_registry, mcp_client, event_bus
    chat_service = ChatService()
    tool_registry = ToolRegistry()
    mcp_client = McpClientManager(tool_registry)
    event_bus = EventBus()

    # Register built-in tools
    shell_tool = ShellTool(mode=os.getenv("TOOLS_MODE", "restricted"))
    fs_tool = FilesystemTool(mode=os.getenv("TOOLS_MODE", "restricted"))
    register_builtin_tools(tool_registry, shell_tool, fs_tool)
    print(f"[mcp] Registered {len(tool_registry.get_all())} built-in tools")

    # Configure Ollama (local LLM)
    ollama_host = os.getenv("OLLAMA_HOST", "http://localhost:11434")
    ollama_model = os.getenv("OLLAMA_MODEL", "llama3.2")
    chat_service.register_provider(OllamaProvider(
        base_url=ollama_host,
        default_model=ollama_model,
    ))

    # Configure Anthropic (if API key is set)
    if os.getenv("ANTHROPIC_API_KEY"):
        chat_service.register_provider(AnthropicProvider())
        print("✓ Anthropic provider enabled")

    # Configure OpenAI (if API key is set)
    if os.getenv("OPENAI_API_KEY"):
        chat_service.register_provider(OpenAIProvider())
        print("✓ OpenAI provider enabled")

    # Configure Google AI / Gemini (if API key is set)
    if os.getenv("GOOGLE_API_KEY"):
        chat_service.register_provider(GoogleProvider())
        print("✓ Google AI (Gemini) provider enabled")

    print(f"🤖 AI Studio Sidecar initialized")
    print(f"   Ollama: {ollama_host} (model: {ollama_model})")

    yield

    # Cleanup: disconnect all MCP servers
    await mcp_client.shutdown()


app = FastAPI(
    title="AI Studio Agent",
    description="Multi-provider LLM agent server with MCP tool support",
    version="0.2.0",
    lifespan=lifespan,
)

# Optional auth token injected by the desktop app.
AI_STUDIO_TOKEN = os.getenv("AI_STUDIO_TOKEN")

@app.middleware("http")
async def auth_middleware(request: Request, call_next):
    if not AI_STUDIO_TOKEN:
        return await call_next(request)

    # Allow unauthenticated health checks (used for startup probing).
    if request.url.path == "/health":
        return await call_next(request)

    token = request.headers.get("x-ai-studio-token")
    if token != AI_STUDIO_TOKEN:
        return JSONResponse(status_code=401, content={"detail": "Unauthorized"})

    return await call_next(request)

# CORS for Tauri/React frontend (used mainly when running UI in a browser).
app.add_middleware(
    CORSMiddleware,
    allow_origins=["tauri://localhost", "http://localhost:1420"] if AI_STUDIO_TOKEN else ["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ============================================================================
# Dynamic Provider Creation
# ============================================================================

def create_provider_for_request(
    name: str,
    api_key: Optional[str] = None,
    base_url: Optional[str] = None,
    extra_config: Optional[dict] = None,
) -> AgentProvider:
    """Create a provider instance from per-request config."""
    cfg = extra_config or {}

    if name == "anthropic":
        return AnthropicProvider(api_key=api_key or "")
    elif name == "google":
        return GoogleProvider(api_key=api_key or "")
    elif name == "azure_openai":
        return AzureOpenAIProvider(
            endpoint=base_url or cfg.get("endpoint", ""),
            api_key=api_key or "",
            deployment=cfg.get("deployment", ""),
            api_version=cfg.get("api_version", "2024-02-01"),
        )
    elif name == "local":
        return LocalOpenAIProvider(
            base_url=base_url or cfg.get("base_url", "http://localhost:11434/v1"),
            api_key=api_key or "",
            model_name=cfg.get("model_name", ""),
        )
    elif name == "openai":
        return OpenAIProvider(api_key=api_key or "")
    elif name == "ollama":
        return OllamaProvider(
            base_url=base_url or "http://localhost:11434",
            default_model=cfg.get("model_name", "llama3.2"),
        )
    else:
        raise ValueError(f"Unknown provider: {name}")


# ============================================================================
# Health & Status Endpoints
# ============================================================================

@app.get("/health")
async def health():
    """Health check endpoint"""
    return {"status": "healthy", "version": "0.2.0"}


@app.get("/status")
async def status():
    """Detailed status including provider health and MCP"""
    provider_health = await chat_service.health_check()
    mcp_servers = {
        name: "connected" for name in mcp_client.get_connected_servers()
    }
    return {
        "status": "healthy",
        "providers": provider_health,
        "mcp_servers": mcp_servers,
        "active_conversations": len(chat_service.conversations),
    }


@app.get("/providers")
async def list_providers():
    """List available LLM providers and their models"""
    return {"providers": chat_service.list_providers()}


# ============================================================================
# Provider Test Endpoint
# ============================================================================

class ProviderTestRequest(BaseModel):
    """Test a provider connection with given credentials"""
    provider: str
    api_key: Optional[str] = None
    base_url: Optional[str] = None
    extra_config: Optional[dict] = None


@app.post("/providers/test")
async def test_provider(request: ProviderTestRequest):
    """Test if a provider connection works with the given config."""
    try:
        provider = create_provider_for_request(
            name=request.provider,
            api_key=request.api_key,
            base_url=request.base_url,
            extra_config=request.extra_config,
        )
        healthy = await provider.health()
        return {"success": healthy, "message": "Connected" if healthy else "Health check failed"}
    except Exception as e:
        return {"success": False, "message": str(e)}


# ============================================================================
# Embedding Endpoint (RAG Knowledge Base)
# ============================================================================

class EmbedRequest(BaseModel):
    """Embedding request for RAG indexing and search"""
    texts: list[str]
    provider: str
    model: str
    # Per-request provider config (from Rust, loaded from SQLite settings)
    api_key: Optional[str] = None
    base_url: Optional[str] = None
    extra_config: Optional[dict] = None


class EmbedResponse(BaseModel):
    """Embedding response"""
    vectors: list[list[float]]
    model: str
    dimensions: int
    usage: dict


@app.post("/embed", response_model=EmbedResponse)
async def embed(request: EmbedRequest):
    """Generate embeddings for a list of texts. Used by Knowledge Base node."""
    try:
        if not request.texts:
            raise HTTPException(status_code=400, detail="texts list cannot be empty")
        if len(request.texts) > 10000:
            raise HTTPException(status_code=400, detail="Maximum 10000 texts per request")

        client = create_embedding_client(
            provider=request.provider,
            api_key=request.api_key or "",
            base_url=request.base_url or "",
            extra_config=request.extra_config,
        )
        result = await client.embed(request.texts, request.model)

        return EmbedResponse(
            vectors=result.vectors,
            model=result.model,
            dimensions=result.dimensions,
            usage=result.usage,
        )
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Embedding error: {str(e)}")


# ============================================================================
# Document Text Extraction (RAG Knowledge Base)
# ============================================================================

class ExtractRequest(BaseModel):
    """Extract text from a binary document (PDF, DOCX, XLSX, PPTX)"""
    path: str
    format: Optional[str] = None  # auto-detected from extension if omitted


class ExtractResponse(BaseModel):
    """Extracted text with metadata"""
    text: str
    format: str
    pages: Optional[int] = None
    sheets: Optional[list[str]] = None
    slides: Optional[int] = None
    char_count: int


def _detect_format(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    format_map = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.xlsx': 'xlsx',
        '.xls': 'xlsx',
        '.pptx': 'pptx',
    }
    fmt = format_map.get(ext)
    if not fmt:
        raise ValueError(f"Unsupported file format: {ext}")
    return fmt


def _extract_pdf(path: str) -> dict:
    from pypdf import PdfReader
    reader = PdfReader(path)
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ''
        if text.strip():
            pages.append(text)
    return {
        'text': '\n\n'.join(pages),
        'pages': len(reader.pages),
    }


def _extract_docx(path: str) -> dict:
    import docx
    doc = docx.Document(path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract table content
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append('\t'.join(cells))
    return {'text': '\n'.join(parts)}


def _extract_xlsx(path: str) -> dict:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    sheet_names = wb.sheetnames
    parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            if any(c for c in cells):
                rows.append('\t'.join(cells))
        if rows:
            parts.append(f'## Sheet: {ws.title}\n' + '\n'.join(rows))
    wb.close()
    return {
        'text': '\n\n'.join(parts),
        'sheets': sheet_names,
    }


def _extract_pptx(path: str) -> dict:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            slides.append(f'## Slide {i}\n' + '\n'.join(texts))
    return {
        'text': '\n\n'.join(slides),
        'slides': len(prs.slides),
    }


@app.post("/extract", response_model=ExtractResponse)
async def extract_text(request: ExtractRequest):
    """Extract text from PDF, DOCX, XLSX, or PPTX files. Used by Knowledge Base node."""
    try:
        if not os.path.exists(request.path):
            raise HTTPException(status_code=404, detail=f"File not found: {request.path}")

        fmt = request.format or _detect_format(request.path)

        extractors = {
            'pdf': _extract_pdf,
            'docx': _extract_docx,
            'xlsx': _extract_xlsx,
            'pptx': _extract_pptx,
        }

        extractor = extractors.get(fmt)
        if not extractor:
            raise HTTPException(status_code=400, detail=f"Unsupported format: {fmt}")

        result = extractor(request.path)
        text = result.get('text', '')

        return ExtractResponse(
            text=text,
            format=fmt,
            pages=result.get('pages'),
            sheets=result.get('sheets'),
            slides=result.get('slides'),
            char_count=len(text),
        )
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Extraction error: {str(e)}")


# ============================================================================
# Chat Endpoints
# ============================================================================

@app.post("/chat", response_model=ChatResponse)
async def chat(request: ChatRequest):
    """
    Send a chat message and get a response.
    Supports tool calling when tools_enabled=True.
    """
    try:
        # Generate conversation ID if not provided
        conversation_id = request.conversation_id or f"conv_{uuid.uuid4().hex[:8]}"

        # If per-request provider config provided, create and register provider on-the-fly
        if request.provider and (request.api_key or request.base_url or request.extra_config):
            provider = create_provider_for_request(
                name=request.provider,
                api_key=request.api_key,
                base_url=request.base_url,
                extra_config=request.extra_config,
            )
            chat_service.register_provider(provider)

        # Create conversation with system prompt if provided
        if request.system_prompt and conversation_id not in chat_service.conversations:
            chat_service.create_conversation(
                conversation_id,
                provider_name=request.provider,
                system_prompt=request.system_prompt,
            )

        # Hydrate from Rust-provided history (SQLite is source of truth)
        if request.history is not None:
            conv = chat_service.get_or_create_conversation(
                conversation_id, provider_name=request.provider,
                system_prompt=request.system_prompt,
            )
            # Build messages: keep system prompt if present, replace the rest.
            # Strip the last user message from history — chat()/chat_with_tools()
            # will re-append it from request.message, avoiding duplication (R1 fix).
            system_msgs = [m for m in conv.messages if m.role == "system"]
            history_msgs = [Message(role=m["role"], content=m["content"]) for m in request.history]
            if history_msgs and history_msgs[-1].role == "user":
                history_msgs = history_msgs[:-1]
            conv.messages = system_msgs + history_msgs

        # Get tool definitions for the provider
        tool_definitions = None
        if request.tools_enabled and tool_registry and len(tool_registry.get_all()) > 0:
            if request.provider == "anthropic":
                tool_definitions = tool_registry.get_anthropic_tools()
            elif request.provider == "google":
                tool_definitions = tool_registry.get_google_tools()
            # Other providers: tools not yet supported, silently skip

        if tool_definitions:
            # Use tool-calling loop
            result = await chat_service.chat_with_tools(
                conversation_id=conversation_id,
                user_message=request.message,
                provider_name=request.provider,
                model=request.model,
                temperature=request.temperature,
                tool_definitions=tool_definitions,
                tool_registry=tool_registry,
                mcp_client=mcp_client,
                event_bus=event_bus,
            )

            # Serialize tool calls for the response
            serialized_tools = [
                {
                    "tool_call_id": tc.tool_call_id,
                    "tool_name": tc.display_name,
                    "tool_input": tc.tool_input,
                    "tool_output": tc.tool_output,
                    "duration_ms": tc.duration_ms,
                    "error": tc.error,
                }
                for tc in result.tool_calls
            ]

            return ChatResponse(
                conversation_id=conversation_id,
                content=result.response.content,
                model=result.response.model,
                provider=result.response.provider,
                usage={
                    "prompt_tokens": result.total_input_tokens,
                    "completion_tokens": result.total_output_tokens,
                },
                tool_calls=serialized_tools if serialized_tools else None,
            )
        else:
            # Simple chat (no tools)
            if event_bus:
                await event_bus.emit(
                    "llm.request.started", conversation_id, "sidecar.chat",
                    {"model": request.model or "", "provider": request.provider or ""},
                )

            llm_start = time.monotonic()
            try:
                response = await chat_service.chat(
                    conversation_id=conversation_id,
                    user_message=request.message,
                    provider_name=request.provider,
                    model=request.model,
                    temperature=request.temperature,
                )
            except Exception as e:
                llm_duration = int((time.monotonic() - llm_start) * 1000)
                if event_bus:
                    await event_bus.emit(
                        "llm.response.error", conversation_id, "sidecar.chat",
                        {
                            "error": str(e),
                            "error_code": type(e).__name__,
                            "model": request.model or "",
                            "provider": request.provider or "",
                            "duration_ms": llm_duration,
                        },
                    )
                raise
            llm_duration = int((time.monotonic() - llm_start) * 1000)

            # Emit llm.response.completed event
            if event_bus:
                await event_bus.emit(
                    "llm.response.completed", conversation_id, "sidecar.chat",
                    {
                        "content": response.content[:500],
                        "model": response.model,
                        "provider": response.provider,
                        "input_tokens": (response.usage or {}).get("prompt_tokens", 0),
                        "output_tokens": (response.usage or {}).get("completion_tokens", 0),
                        "duration_ms": llm_duration,
                        "stop_reason": response.stop_reason or "end_turn",
                    },
                )

            return ChatResponse(
                conversation_id=conversation_id,
                content=response.content,
                model=response.model,
                provider=response.provider,
                usage=response.usage,
            )

    except ValueError as e:
        if event_bus and conversation_id:
            await event_bus.emit(
                "agent.error", conversation_id, "sidecar.chat",
                {"error": str(e), "error_code": "ValueError", "severity": "warning"},
            )
        raise HTTPException(status_code=400, detail=str(e))
    except HTTPException:
        raise
    except Exception as e:
        if event_bus and conversation_id:
            await event_bus.emit(
                "agent.error", conversation_id, "sidecar.chat",
                {"error": str(e), "error_code": type(e).__name__, "severity": "error"},
            )
        raise HTTPException(status_code=500, detail=f"Chat error: {str(e)}")


def _prepare_chat_request(request: ChatMessageRequest):
    """Shared setup for /chat/direct and /chat/stream: provider, messages, session."""
    import json as _json
    print(f"[chat] provider={request.provider} model={request.model} "
          f"base_url={request.base_url} extra_config={request.extra_config} "
          f"msgs={len(request.messages)} session={request.conversation_id or 'none'}")

    # Register provider on-the-fly if config provided
    if request.provider and (request.api_key or request.base_url or request.extra_config):
        provider_inst = create_provider_for_request(
            name=request.provider,
            api_key=request.api_key,
            base_url=request.base_url,
            extra_config=request.extra_config,
        )
        chat_service.register_provider(provider_inst)
    else:
        print(f"[chat] WARN: No dynamic config — using default provider for '{request.provider}'")

    messages = []
    if request.system_prompt:
        messages.append(Message(role="system", content=request.system_prompt))

    # Session mode: inject history from previous turns
    conv = None
    if request.conversation_id:
        conv = chat_service.get_or_create_conversation(
            request.conversation_id,
            provider_name=request.provider,
            system_prompt=request.system_prompt,
        )
        history = [m for m in conv.messages if m.role != "system"]
        max_h = request.max_history or 20
        if len(history) > max_h:
            history = history[-max_h:]
        if history:
            print(f"[chat] Session '{request.conversation_id}': injecting {len(history)} history messages")
            messages.extend(history)

    # Build current messages — inject images into multimodal content if present
    if request.images:
        print(f"[chat] Vision mode: {len(request.images)} image(s) attached")
        for m in request.messages:
            content_blocks = [{"type": "text", "text": m["content"]}]
            for img in request.images:
                data_uri = f"data:{img['mime_type']};base64,{img['data']}"
                content_blocks.append({
                    "type": "image_url",
                    "image_url": {"url": data_uri},
                })
            messages.append(Message(role=m["role"], content=content_blocks))
    else:
        messages.extend(Message(role=m["role"], content=m["content"]) for m in request.messages)

    provider = chat_service.get_provider(request.provider or "ollama")
    return provider, messages, conv


@app.post("/chat/direct")
async def chat_direct(request: ChatMessageRequest):
    """Direct chat without conversation history. Used by workflow LLM nodes.
    When conversation_id is provided, accumulates history across calls (session mode).
    """
    try:
        provider, messages, conv = _prepare_chat_request(request)

        response = await provider.chat(
            messages=messages,
            model=request.model,
            temperature=request.temperature,
        )

        # Session mode: store current turn in conversation history
        if conv is not None:
            for m in request.messages:
                conv.messages.append(Message(role=m["role"], content=m["content"]))
            conv.messages.append(Message(role="assistant", content=response.content))

        return {
            "content": response.content,
            "model": response.model,
            "provider": response.provider,
            "usage": response.usage,
        }

    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Chat error: {str(e)}")


@app.post("/chat/stream")
async def chat_stream_endpoint(request: ChatMessageRequest):
    """Streaming chat via SSE. Same request schema as /chat/direct.
    Returns text/event-stream with token/done/error chunks.
    """
    import json as _json

    try:
        provider, messages, conv = _prepare_chat_request(request)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Chat setup error: {str(e)}")

    async def generate():
        accumulated = ""
        try:
            async for chunk in provider.chat_stream(
                messages=messages,
                model=request.model,
                temperature=request.temperature,
            ):
                if chunk['type'] == 'token':
                    accumulated += chunk['content']
                elif chunk['type'] == 'done':
                    accumulated = chunk.get('content', accumulated)
                    # Session mode: store history after stream completes
                    if conv is not None:
                        for m in request.messages:
                            conv.messages.append(Message(role=m["role"], content=m["content"]))
                        conv.messages.append(Message(role="assistant", content=accumulated))
                yield f"data: {_json.dumps(chunk)}\n\n"
        except Exception as e:
            print(f"[chat/stream] ERROR: {e}")
            yield f"data: {_json.dumps({'type': 'error', 'message': str(e)})}\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


class ToolExecuteRequest(BaseModel):
    """Execute a single tool by name"""
    tool_name: str
    tool_input: dict = {}


@app.post("/tools/execute")
async def execute_tool(request: ToolExecuteRequest):
    """Execute a tool directly by name (used by workflow Tool nodes)."""
    try:
        tool_def = tool_registry.resolve(request.tool_name)
        if tool_def and tool_def.handler:
            result = await tool_def.handler(**request.tool_input)
            return {"result": result}

        # Try MCP servers
        parts = request.tool_name.split("__", 1)
        if len(parts) == 2:
            server, local_name = parts
            if mcp_client.is_connected(server):
                result = await mcp_client.call_tool(server, local_name, request.tool_input)
                return {"result": result}

        raise HTTPException(status_code=404, detail=f"Tool '{request.tool_name}' not found")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Tool execution error: {str(e)}")


@app.delete("/chat/{conversation_id}")
async def delete_conversation(conversation_id: str):
    """Delete a conversation"""
    chat_service.delete_conversation(conversation_id)
    return {"status": "deleted", "conversation_id": conversation_id}


@app.post("/chat/{conversation_id}/clear")
async def clear_conversation(conversation_id: str):
    """Clear a conversation's history"""
    chat_service.clear_conversation(conversation_id)
    return {"status": "cleared", "conversation_id": conversation_id}


# ============================================================================
# MCP Endpoints
# ============================================================================

@app.post("/mcp/connect")
async def mcp_connect(request: McpConnectRequest):
    """Connect to an MCP server and discover its tools."""
    from agent.mcp.client import McpServerConfig
    config = McpServerConfig(
        name=request.name,
        transport=request.transport,
        command=request.command,
        args=request.args,
        url=request.url,
        env=request.env,
    )
    result = await mcp_client.connect(config)
    return result


@app.post("/mcp/disconnect")
async def mcp_disconnect(request: McpDisconnectRequest):
    """Disconnect from an MCP server."""
    await mcp_client.disconnect(request.name)
    return {"status": "disconnected", "name": request.name}


@app.get("/mcp/tools")
async def mcp_tools():
    """List all available tools (built-in + MCP)."""
    return {"tools": tool_registry.to_summary()}


# ============================================================================
# WebSocket — Live Event Stream
# ============================================================================

@app.websocket("/events")
async def events_ws(websocket: WebSocket):
    """
    Live event stream for the Inspector.

    Auth: first message must be {"type": "auth", "token": "<AI_STUDIO_TOKEN>"}.
    After auth, events are streamed as JSON (one per WS message).
    """
    await websocket.accept()

    # Auth handshake
    if AI_STUDIO_TOKEN:
        try:
            msg = await asyncio.wait_for(websocket.receive_json(), timeout=5.0)
            if msg.get("type") != "auth" or msg.get("token") != AI_STUDIO_TOKEN:
                await websocket.close(code=4001, reason="Unauthorized")
                return
        except Exception:
            await websocket.close(code=4001, reason="Auth timeout")
            return

    queue = event_bus.subscribe()
    print("[ws] Event subscriber connected")
    try:
        while True:
            msg = await queue.get()
            await websocket.send_text(msg)
    except WebSocketDisconnect:
        pass
    except Exception:
        pass
    finally:
        event_bus.unsubscribe(queue)
        print("[ws] Event subscriber disconnected")


# ============================================================================
# Tools Endpoints (Legacy — direct tool access)
# ============================================================================

# Initialize tools with restricted mode by default
shell_tool = ShellTool(mode=os.getenv("TOOLS_MODE", "restricted"))
filesystem_tool = FilesystemTool(mode=os.getenv("TOOLS_MODE", "restricted"))

from agent.tools import BrowserTool
browser_tool = BrowserTool(headless=True)


class ShellRequest(BaseModel):
    """Shell command request"""
    command: str
    timeout: Optional[float] = 30.0
    cwd: Optional[str] = None


class FilesystemRequest(BaseModel):
    """Filesystem operation request"""
    action: str  # read, write, delete, list_dir, mkdir, exists, copy, move
    path: str
    content: Optional[str] = None  # For write/append
    dest: Optional[str] = None  # For copy/move


class BrowserRequest(BaseModel):
    """Browser action request"""
    action: str  # navigate, screenshot, click, fill, extract_text, get_html, evaluate
    url: Optional[str] = None
    selector: Optional[str] = None
    value: Optional[str] = None
    script: Optional[str] = None


@app.post("/tools/shell")
async def run_shell(request: ShellRequest):
    """Execute a shell command."""
    result = await shell_tool.run(
        command=request.command,
        timeout=request.timeout,
        cwd=request.cwd,
    )
    return {
        "command": result.command,
        "stdout": result.stdout,
        "stderr": result.stderr,
        "return_code": result.return_code,
        "timed_out": result.timed_out,
    }


@app.post("/tools/filesystem")
async def run_filesystem(request: FilesystemRequest):
    """Perform filesystem operations."""
    action = request.action.lower()

    if action == "read":
        result = filesystem_tool.read(request.path)
    elif action == "write":
        result = filesystem_tool.write(request.path, request.content or "")
    elif action == "append":
        result = filesystem_tool.append(request.path, request.content or "")
    elif action == "delete":
        result = filesystem_tool.delete(request.path)
    elif action == "list_dir":
        result = filesystem_tool.list_dir(request.path)
    elif action == "mkdir":
        result = filesystem_tool.mkdir(request.path)
    elif action == "exists":
        result = filesystem_tool.exists(request.path)
    elif action == "copy":
        result = filesystem_tool.copy(request.path, request.dest or "")
    elif action == "move":
        result = filesystem_tool.move(request.path, request.dest or "")
    else:
        raise HTTPException(status_code=400, detail=f"Unknown action: {action}")

    return {
        "success": result.success,
        "action": result.action,
        "path": result.path,
        "data": result.data,
        "error": result.error,
    }


@app.post("/tools/browser/start")
async def browser_start():
    """Start the browser instance"""
    result = await browser_tool.start()
    return {"success": result.success, "error": result.error}


@app.post("/tools/browser/stop")
async def browser_stop():
    """Stop the browser instance"""
    result = await browser_tool.stop()
    return {"success": result.success, "error": result.error}


@app.post("/tools/browser")
async def run_browser(request: BrowserRequest):
    """Perform browser actions."""
    action = request.action.lower()

    if action == "navigate":
        result = await browser_tool.navigate(request.url or "")
    elif action == "screenshot":
        result = await browser_tool.screenshot()
    elif action == "click":
        result = await browser_tool.click(request.selector or "")
    elif action == "fill":
        result = await browser_tool.fill(request.selector or "", request.value or "")
    elif action == "extract_text":
        result = await browser_tool.extract_text(request.selector or "body")
    elif action == "get_html":
        result = await browser_tool.get_html(request.selector or "body")
    elif action == "evaluate":
        result = await browser_tool.evaluate(request.script or "")
    elif action == "wait_for":
        result = await browser_tool.wait_for(request.selector or "")
    else:
        raise HTTPException(status_code=400, detail=f"Unknown action: {action}")

    return {
        "success": result.success,
        "action": result.action,
        "data": result.data,
        "screenshot": result.screenshot,
        "error": result.error,
    }


# ============================================================================
# Main Entry Point
# ============================================================================

if __name__ == "__main__":
    import uvicorn

    host = os.getenv("HOST", "127.0.0.1")
    port = int(os.getenv("PORT", "8765"))

    print(f"\n🚀 Starting AI Studio Agent Server")
    print(f"   URL: http://{host}:{port}")
    print(f"   Docs: http://{host}:{port}/docs")
    print(f"\n   Press Ctrl+C to stop\n")

    uvicorn.run(app, host=host, port=port)
